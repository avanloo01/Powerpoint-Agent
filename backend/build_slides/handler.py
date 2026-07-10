"""
Build-slides Lambda: generates python-pptx code from a presentation structure,
executes it, and uploads the resulting PPTX to S3.
Invoked asynchronously by agent_loop after research & structure are complete.
"""
from __future__ import annotations

import builtins as _builtins
import io
import json
import math
import os
import re
import textwrap
import time
import traceback
from urllib import request as urlrequest

import boto3
from openai import OpenAI

# ─── CONSTANTS ────────────────────────────────────────────────────────────────

S3_OUTPUT_BUCKET = os.environ.get("OUTPUT_BUCKET", "")
QWEN_BASE_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1"
QWEN_MODEL = os.environ.get("QWEN_MODEL", "qwen3.6-plus")
SUPABASE_URL = os.environ.get("SUPABASE_URL", "").rstrip("/")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")

# ─── SUPABASE HELPERS ─────────────────────────────────────────────────────────

def _supabase_request(method: str, url: str, headers: dict, body: dict | None = None) -> object:
    payload = None
    request_headers = dict(headers)
    if body is not None:
        payload = json.dumps(body).encode("utf-8")
        request_headers["Content-Type"] = "application/json"
    req = urlrequest.Request(url=url, method=method, data=payload, headers=request_headers)
    try:
        with urlrequest.urlopen(req, timeout=10) as resp:
            raw = resp.read().decode("utf-8")
            return json.loads(raw) if raw else {}
    except Exception:  # noqa: BLE001
        return {}


def _update_job(job_id: str, **fields: object) -> None:
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        return
    url = f"{SUPABASE_URL}/rest/v1/jobs?id=eq.{job_id}"
    _supabase_request(
        "PATCH",
        url,
        {
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
            "Prefer": "return=minimal",
        },
        fields,  # type: ignore[arg-type]
    )


# ─── ICON RENDERING (Pillow-only, no cairo needed) ────────────────────────────

def _recolor_png(png_bytes: bytes, hex_color: str) -> bytes:
    """Recolor a single-colour PNG icon by replacing the RGB of all
    non-transparent pixels with the target hex colour, preserving alpha."""
    from io import BytesIO
    from PIL import Image

    img = Image.open(BytesIO(png_bytes)).convert("RGBA")
    data = bytearray(img.tobytes())
    w, h = img.size
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    for y in range(h):
        for x in range(w):
            idx = (y * w + x) * 4
            if data[idx + 3] > 0:  # non-transparent pixel
                data[idx] = r
                data[idx + 1] = g
                data[idx + 2] = b
    out = BytesIO()
    Image.frombytes("RGBA", (w, h), bytes(data)).save(out, format="PNG")
    out.seek(0)
    return out.read()


# ─── BUILD AGENT ─────────────────────────────────────────────────────────────

_BUILD_SYSTEM = textwrap.dedent("""\
You are an expert python-pptx developer. Write a Python 3.12 function called `build_presentation(prs, logo_bytes=None)` that adds all slides to the given python-pptx Presentation object `prs`. ONLY return the function code, no markdown.

CONSTRAINTS:
- prs: slide_width=Inches(13.33), slide_height=Inches(7.5). Use slide_layouts[6] (blank).
- Never call Presentation(). Never write import statements. Pre-injected globals:
  Inches, Pt, Emu, Cm, RGBColor, PP_ALIGN, ChartData, XL_CHART_TYPE, MSO_SHAPE,
  MSO_ANCHOR, MSO_AUTO_SIZE, io, math, json, BytesIO, Image, ImageEnhance,
  urlrequest, no_shadow, remove_outline, get_image_buf
- slides_data is injected as a FUNCTION PARAMETER (list of slide dicts). Iterate:
  `for slide in slides_data:` — NEVER re-define slides_data or embed slide data as Python dict/list literals.
- Always use enum constants, not integers: MSO_ANCHOR.MIDDLE not 2; PP_ALIGN.CENTER not 1.

CRITICAL RULES (violating these WILL crash):
- add_picture() needs a file-like object, not raw bytes:
  buf.seek(0); slide.shapes.add_picture(buf, ...)  ← always seek(0) after any BytesIO write.
- Images: NEVER define fetch_bg/get_bg/download_image or call urlrequest for images.
  Only pattern: buf = get_image_buf(url, darken=True/False); slide.shapes.add_picture(buf, ...)
  get_image_buf() is pre-injected, returns a seek'd BytesIO, handles failures gracefully.
- no_shadow(shape): call on every add_shape result. Never shape.shadow.inherit = False.
- remove_outline(shape): call on every filled shape to strip the default blue border.
  Never shape.line.fill.background() directly — crashes on charts/tables/GraphicFrame.

STYLE GUIDE:
- COLORS: PRIMARY = RGBColor(...) from prompt for all accents. ACCENT for chart series only.
    Also capture your primary R,G,B as: R1,G1,B1 = 0xC0,0x00,0x00 (the same values you used for PRIMARY).
    These are used by the pie-chart shading template to create progressive lighter shades.
- LAYOUT CONSTANTS (define once at top of each content-slide block):
    CONCLUSION_Y   = sh - Cm(2.4)
    CONCLUSION_H   = Cm(1.2)
    SOURCES_Y      = sh - Cm(0.9)
    CONTENT_BOTTOM = CONCLUSION_Y - Cm(0.2)
- title_slide: if image_url provided: get_image_buf(url, darken=True) → add_picture(buf,0,0,sw,sh).
  If image_url is null: add_shape(RECTANGLE,0,0,sw,sh) filled RGBColor(40,40,40), no_shadow, remove_outline.
  Title 54pt bold white centered.
- section_divider: if image_url provided: get_image_buf(url, darken=True) → add_picture(buf,0,0,sw,sh).
  If image_url is null: add_shape(RECTANGLE,0,0,sw,sh) filled RGBColor(40,40,40), no_shadow, remove_outline.
  Then:
    wy = sh - Cm(5.74)
    a) White rect (0, wy, sw, Cm(5.74))
    b) Number square (0, wy, Cm(5.74), Cm(5.74)): PRIMARY fill, "{n:02d}" white bold 48pt, CENTER+MIDDLE
    c) Title textbox (Cm(6.27), wy+Cm(0.8), sw-Cm(7.27), Cm(2.5)): bold black 32pt, word_wrap=True
       tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  ← required so .height is accurate
    d) Section label: y = title_tb.top + title_tb.height + Cm(0.2), gray 16pt. NEVER hardcode y.
- content_slide: ALWAYS set `slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255,255,255)` as the FIRST two lines of every content-slide block — skipping this leaves a grey default background. No white rect shape.
  FIXED positions (do NOT derive from margin — they must clear BOX_HEADER_Y=Cm(3.05)):
    Section label: x=margin, y=Cm(0.5), w=sw-2*margin, h=Cm(0.5), 9pt gray, word_wrap=True
    Title:         x=margin, y=Cm(1.1), w=sw-2*margin, h=Cm(1.8), 22pt bold black, word_wrap=True
- box_headers: FIXED coordinates — define once at the top of each content-slide block:
    BOX_HEADER_Y = Cm(3.05)
    BOX_HEADER_H = Cm(0.81)
  PRIMARY rect at (col_x+Inches(0.08), BOX_HEADER_Y, col_w-Inches(0.16), BOX_HEADER_H),
  white bold 12pt, margin_left=Inches(0.1).
- column subtitle: if a column's "subtitle" field is non-null, add a textbox
  (col_x+Inches(0.08), BOX_HEADER_Y+BOX_HEADER_H, col_w-Inches(0.16), Cm(0.4)): 8pt gray italic text.
  Increment content start y by Cm(0.4) + Inches(0.05) before rendering chart/bullets.
- bullets: icons[filename] → BytesIO → add_picture Pt(22)×Pt(22). Title bold 10pt, description 8pt.
  ALWAYS start with `bullet_top = content_y + Cm(0.2)` — this mandatory top gap prevents bullets from crowding the header box.
  Icon y-centred with title+desc block: icon_y = bullet_top + (title_h+desc_h+gap - Pt(22))/2.
  BULLET SPACING: after each bullet, advance bullet_top by title_h + desc_h + gap + Pt(10). The Pt(10) inter-bullet padding is mandatory — never use a flat Cm(1.0) increment that ignores it.
  Icon fallback: if the icons dict is empty, use MSO_SHAPE.OVAL Pt(10) PRIMARY fill.
  CRITICAL: NEVER render Unicode symbols (↑ ↓ → ← ✓ ✗ ● ◆ •) as text characters in title or description — the icon field is the only place for visual markers.
- separators: usable = sw - 2*margin - gap; col_w = int(usable * width_ratio).  ← MUST use int() — float EMU values corrupt the PPTX XML.
    LIGHT_GRAY = RGBColor(211, 211, 211)
    sep_x = col1_x + col1_w + gap/2 - Pt(0.25)
    line_top = min(col_tops); line_bottom = CONTENT_BOTTOM
    Draw Pt(0.5) × (line_bottom - line_top) LIGHT_GRAY rect at sep_x.
    "line": rect only. No triangle.
    "causal_line": rect + right-pointing triangle whose base sits on the separator line:
      # After rotation=90 the visual base is at tri_cx-Pt(4); shift tri_cx right so base lands on the line.
      tri_cx = sep_x + Pt(4.25)  # base at tri_cx-Pt(4) = sep_x+Pt(0.25) ← on the line
      tri_cy = (line_top + line_bottom) / 2
      tri = add_shape(ISOSCELES_TRIANGLE, tri_cx-Pt(6), tri_cy-Pt(4), Pt(12), Pt(8))
      tri.rotation = 90  # apex points right; visual width=Pt(8), visual height=Pt(12)
      tri.fill.solid(); tri.fill.fore_color.rgb = LIGHT_GRAY; no_shadow(tri); remove_outline(tri)
- conclusion_box: if non-null, draw BEFORE sources using MSO_SHAPE.RECTANGLE (NEVER ROUNDED_RECTANGLE):
    box = add_shape(RECTANGLE, margin, CONCLUSION_Y, sw-2*margin, CONCLUSION_H)
    box.fill.background(); no_shadow(box)  ← do NOT remove_outline — border is intentional
    box.line.color.rgb = PRIMARY; box.line.width = Pt(1.0)
    tf: word_wrap=True, MIDDLE anchor, 0.15in margins, 9pt black PP_ALIGN.CENTER text
- sources: draw LAST at SOURCES_Y. 8pt gray left-aligned.
  RENDER ORDER: columns → separator → conclusion_box → sources → logo
- logo: top-right, ~0.6in tall, BytesIO(logo_bytes).
- charts: Use EXACTLY this pattern — all four types are handled via a dict map with fallback:
    _CHART_TYPE_MAP = {
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'grouped_bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
    }
    ct = _CHART_TYPE_MAP.get(ch['chart_type'], XL_CHART_TYPE.COLUMN_CLUSTERED)
    cd = ChartData()
    cd.categories = ch['x_labels']
    for s in ch['series']:
        cd.add_series(s['name'], s['values'])
    content_y = BOX_HEADER_Y + BOX_HEADER_H + Inches(0.1)
    chart_h = max(CONTENT_BOTTOM - content_y - Inches(0.1), Inches(2.5))
    # if column also has bullets below chart: chart_h = (CONTENT_BOTTOM - content_y) * 0.55
    cf = slide.shapes.add_chart(ct, col_x+Inches(0.08), content_y, col_w-Inches(0.16), chart_h, cd)
    chart_obj = cf.chart
    chart_obj.chart_style = 2
    if ch['chart_type'] in ('bar', 'grouped_bar'):
        chart_obj.value_axis.has_major_gridlines = False
        chart_obj.value_axis.visible = False
        plot = chart_obj.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True
        for i, ser in enumerate(chart_obj.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = PRIMARY if i == 0 else ACCENT
    elif ch['chart_type'] == 'line':
        chart_obj.value_axis.has_major_gridlines = False
        for i, ser in enumerate(chart_obj.series):
            ser.format.line.color.rgb = PRIMARY if i == 0 else ACCENT
    elif ch['chart_type'] == 'pie':
        plot = chart_obj.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_category_name = True
        plot.data_labels.show_percentage = True
        # Add a legend at the bottom so slice categories are visible
        chart_obj.has_legend = True
        chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart_obj.legend.include_in_layout = False
        # Color slices: progressive lighter shades of PRIMARY (blend toward white).
        # Use the numeric RGB components from your PRIMARY definition (e.g. R1=0xC0, G1=0x00, B1=0x00).
        for j, pt in enumerate(chart_obj.series[0].points):
            t = (j + 1) / (len(chart_obj.series[0].points) + 1)
            shade = RGBColor(
                int(R1 + (255 - R1) * t),
                int(G1 + (255 - G1) * t),
                int(B1 + (255 - B1) * t),
            )
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = shade
  Never call chart.replace_data() — build all data via ChartData at creation time.
""")


BATCH_SIZE = 4


def _build_batch_code(
    batch_slides: list,
    batch_num: int,
    total_batches: int,
    settings: dict,
    client: OpenAI,
    job_id: str,
    has_logo: bool = False,
    has_icons: bool = False,
) -> str:
    """Generate python-pptx code for a single batch of slides."""
    primary = settings.get("primary_color", "#C00000")
    accent = settings.get("accent_color", "#A6CAEC")
    fn_name = f"build_batch_{batch_num}"
    logo_note = (
        "A `logo_bytes` PARAMETER will be passed to your function — include it in your signature. "
        "Place the logo on every slide."
        if has_logo else "No logo will be provided — your function does NOT need a `logo_bytes` parameter."
    )
    icon_note = (
        "An `icons` dict (filename → PNG bytes) is passed as a PARAMETER to your function — "
        f"include it in your signature: def {fn_name}(prs, logo_bytes=None, icons=None, slides_data=None). "
        "Use icons[bullet['icon']] as a BytesIO source for add_picture()."
        if has_icons else "No icons are provided — omit the `icons` parameter. Use MSO_SHAPE.OVAL Pt(10), PRIMARY fill as fallback."
    )
    batch_structure: dict = {"slides": batch_slides}
    structure_json = json.dumps(batch_structure, indent=2)
    user_msg = (
        f"Primary color: {primary}\n"
        f"Accent color: {accent}\n"
        f"Logo: {logo_note}\n"
        f"Icons: {icon_note}\n"
        f"IMPORTANT: Name your function `{fn_name}`. Your function signature MUST include `slides_data=None`: "
        f"`def {fn_name}(prs, logo_bytes=None, icons=None, slides_data=None):` — omit logo_bytes/icons only if told they're not provided, but ALWAYS include slides_data.\n"
        f"The batch slides are injected as `slides_data` (a Python list). Iterate: `for slide in slides_data:`. "
        f"DO NOT re-define slides_data or encode any slide data as Python dict/list literals in your code body.\n\n"
        f"Presentation structure (reference only — iterate slides_data at runtime, do not copy this JSON into code):\n{structure_json}"
    )
    prompt_len = len(_BUILD_SYSTEM) + len(user_msg)
    last_exc = None
    for retry in range(3):
        try:
            print(f"[{job_id}] Batch {batch_num}/{total_batches}: API call ({prompt_len} chars, timeout=600s, attempt {retry+1}/3)...")
            response = client.chat.completions.create(
                model=QWEN_MODEL,
                messages=[
                    {"role": "system", "content": _BUILD_SYSTEM},
                    {"role": "user", "content": user_msg},
                ],
                max_tokens=16000,
                timeout=600.0,
            )
            code = response.choices[0].message.content or ""
            print(f"[{job_id}] Batch {batch_num}/{total_batches}: received {len(code)} chars")
            return code
        except Exception as exc:
            last_exc = exc
            print(f"[{job_id}] Batch {batch_num}/{total_batches}: API attempt {retry+1} failed: {exc}")
            if retry < 2:
                wait = (retry + 1) * 30
                print(f"[{job_id}] Waiting {wait}s before retry...")
                time.sleep(wait)
    raise last_exc  # type: ignore[misc]


# ─── EXECUTE GENERATED CODE ───────────────────────────────────────────────────

def _make_namespace(image_buffers: dict[str, bytes] | None = None) -> dict:
    """Create the restricted execution namespace with pre-injected globals."""
    import pptx
    import pptx.chart.data
    import pptx.dml.color
    import pptx.enum.chart
    import pptx.enum.text
    import pptx.enum.shapes
    import pptx.util
    from io import BytesIO
    from PIL import Image, ImageEnhance

    # ── Patch pptx to accept float coordinates (AI often divides EMU ints) ──
    # Charts use CT_GraphicalObjectFrame.new_graphicFrame which embeds coords via
    # f-string (f'cx="{cx}"'), bypassing all type validators. Patch it directly.
    import pptx.oxml.shapes.graphfrm as _graphfrm
    import pptx.oxml.simpletypes as _simpletypes

    _orig_new_graphicFrame = _graphfrm.CT_GraphicalObjectFrame.new_graphicFrame

    @classmethod
    def _int_new_graphicFrame(cls, id_: int, name: str, x: int, y: int, cx: int, cy: int):
        return _orig_new_graphicFrame.__func__(cls, id_, name, int(x), int(y), int(cx), int(cy))

    _graphfrm.CT_GraphicalObjectFrame.new_graphicFrame = _int_new_graphicFrame

    # Also patch BaseIntType.convert_to_xml for any remaining descriptor-based setters
    _orig_convert_to_xml = _simpletypes.BaseIntType.convert_to_xml

    @classmethod
    def _lenient_convert_to_xml(cls, value):
        if isinstance(value, float):
            value = int(value)
        return _orig_convert_to_xml.__func__(cls, value)

    _simpletypes.BaseIntType.convert_to_xml = _lenient_convert_to_xml

    # And validate_int so floats don't raise before reaching convert_to_xml
    _orig_validate_int = _simpletypes.BaseSimpleType.validate_int

    @staticmethod
    def _lenient_validate_int(value: int | float) -> None:
        if isinstance(value, float):
            value = int(value)
        return _orig_validate_int(value)

    _simpletypes.BaseSimpleType.validate_int = _lenient_validate_int

    def _no_shadow(shape):
        """Safely disable shadows — silently skips shapes that don't support .shadow (pictures, charts, tables)."""
        try:
            shape.shadow.inherit = False
        except (NotImplementedError, AttributeError):
            pass

    def _remove_outline(shape):
        """Safely remove shape outline — silently skips shapes that don't support .line (charts, tables, GraphicFrame)."""
        try:
            shape.line.fill.background()
        except (NotImplementedError, AttributeError):
            pass

    _SAFE_NAMES = (
        "abs", "all", "any", "bool", "bytes", "bytearray", "dict", "enumerate",
        "filter", "float", "globals", "hasattr", "int",
        "isinstance", "iter", "len", "list", "map", "max", "min", "next", "print", "range",
        "reversed", "round", "set", "sorted", "str", "sum", "tuple", "zip",
        "Exception", "ValueError", "TypeError", "KeyError",
        "IndexError", "AttributeError", "RuntimeError", "StopIteration",
    )
    safe_builtins = {
        name: getattr(_builtins, name)
        for name in _SAFE_NAMES
        if hasattr(_builtins, name)
    }

    ns: dict = {
        "__builtins__": safe_builtins,
        "__image_buffers__": image_buffers or {},
        "Inches": pptx.util.Inches,
        "Pt": pptx.util.Pt,
        "Emu": pptx.util.Emu,
        "Cm": pptx.util.Cm,
        "RGBColor": pptx.dml.color.RGBColor,
        "PP_ALIGN": pptx.enum.text.PP_ALIGN,
        "XL_LEGEND_POSITION": pptx.enum.chart.XL_LEGEND_POSITION,
        "ChartData": pptx.chart.data.ChartData,
        "XL_CHART_TYPE": pptx.enum.chart.XL_CHART_TYPE,
        "MSO_SHAPE": pptx.enum.shapes.MSO_SHAPE,
        "MSO_ANCHOR": pptx.enum.text.MSO_ANCHOR,
        "MSO_AUTO_SIZE": pptx.enum.text.MSO_AUTO_SIZE,
        "io": io,
        "math": math,
        "json": json,
        "urlrequest": urlrequest,
        "Image": Image,
        "ImageEnhance": ImageEnhance,
        "BytesIO": BytesIO,
        "no_shadow": _no_shadow,
        "remove_outline": _remove_outline,
    }
    if image_buffers:
        ns["get_image_buf"] = _get_image_buf_factory(image_buffers)
    return ns


def _get_image_buf_factory(image_buffers: dict[str, bytes]):
    """Return a get_image_buf(url, darken=False) closure that looks up
    pre-downloaded images, wrapping them in BytesIO and optionally darkening."""
    from io import BytesIO
    from PIL import Image, ImageEnhance

    def _get_image_buf(url: str, darken: bool = False):
        data = image_buffers.get(url)
        if data is None:
            # Image failed to download — return a dark placeholder so the slide renders
            placeholder = Image.new("RGB", (1920, 1080), (50, 50, 50))
            buf = BytesIO()
            placeholder.save(buf, format="PNG")
            buf.seek(0)
            return buf
        img = Image.open(BytesIO(data))
        if darken:
            enhancer = ImageEnhance.Brightness(img)
            img = enhancer.enhance(0.6)
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf

    return _get_image_buf


def _rollback_slides(prs: object, original_count: int) -> None:
    """Remove slides added beyond original_count to undo a failed batch execution.
    This prevents orphan slides (often with grey background) from persisting in the PPTX."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    removed = 0
    while len(sldIdLst) > original_count:
        sld_id = sldIdLst[-1]
        r_id = sld_id.get(qn('r:id'))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:  # noqa: BLE001
                pass
        sldIdLst.remove(sld_id)
        removed += 1
    return removed


class _SafeIconDict(dict):
    """dict subclass that returns the first available icon bytes when a key is
    missing, preventing the OVAL shape fallback from triggering on hallucinated
    or mis-cased icon filenames generated by the LLM."""

    def __missing__(self, key: str) -> bytes:
        if self:
            return next(iter(self.values()))
        raise KeyError(key)


def _execute_batch(
    code: str, prs: object, namespace: dict, logo_bytes: bytes | None, fn_name: str,
    icons: dict[str, bytes] | None = None,
    slides_data: list | None = None,
) -> None:
    """Execute AI-generated batch code, adding slides to the existing prs."""
    code = re.sub(r"```(?:python)?\s*", "", code).strip().rstrip("`").strip()
    exec(code, namespace)  # noqa: S102

    build_fn = namespace.get(fn_name)
    if not callable(build_fn):
        available = [k for k, v in namespace.items() if callable(v) and not k.startswith("_")]
        raise ValueError(
            f"Generated code does not define '{fn_name}'. "
            f"Available callables: {available}"
        )

    kwargs = {"logo_bytes": logo_bytes, "icons": _SafeIconDict(icons or {})}
    import inspect
    sig = inspect.signature(build_fn)
    if "image_buffers" in sig.parameters:
        kwargs["image_buffers"] = namespace.get("__image_buffers__", {})
    if "slides_data" in sig.parameters:
        kwargs["slides_data"] = slides_data or []
    build_fn(prs, **kwargs)

# ─── REFERENCES SLIDE ─────────────────────────────────────────────────────
def _add_references_slide(prs: object, structure: dict, settings: dict) -> None:
    """Append a References slide listing all source URLs from slide notes."""
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    slides_data = structure.get("slides", [])
    refs: list[tuple[str, str, str]] = []
    for slide in slides_data:
        notes = (slide.get("notes") or "").strip()
        if not notes:
            continue
        refs.append((
            slide.get("slide_title", ""),
            (slide.get("sources") or "").strip(),
            notes,
        ))

    if not refs:
        return

    sw = prs.slide_width
    sh = prs.slide_height
    ref_slide = prs.slides.add_slide(prs.slide_layouts[6])

    ref_slide.background.fill.solid()
    ref_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    primary_hex = settings.get("primary_color", "#C00000").lstrip("#")
    primary_rgb = RGBColor(
        int(primary_hex[0:2], 16),
        int(primary_hex[2:4], 16),
        int(primary_hex[4:6], 16),
    )

    # Title
    title_tb = ref_slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Cm(1.0))
    p = title_tb.text_frame.paragraphs[0]
    p.text = "References"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = primary_rgb

    # Thin separator bar
    sep = ref_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.0), sw - Inches(0.8), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = primary_rgb
    try:
        sep.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass
    try:
        sep.line.fill.background()
    except Exception:  # noqa: BLE001
        pass

    # References content
    content_tb = ref_slide.shapes.add_textbox(Inches(0.4), Inches(1.15), sw - Inches(0.8), sh - Inches(1.45))
    ctf = content_tb.text_frame
    ctf.word_wrap = True

    first = True
    for title_str, sources_str, notes_str in refs:
        if first:
            p1 = ctf.paragraphs[0]
            first = False
        else:
            p_gap = ctf.add_paragraph()
            p_gap.font.size = Pt(4)
            p1 = ctf.add_paragraph()

        p1.text = title_str + (f"  |  {sources_str}" if sources_str else "")
        p1.font.bold = True
        p1.font.size = Pt(9)
        p1.font.color.rgb = RGBColor(40, 40, 40)

        p2 = ctf.add_paragraph()
        p2.text = notes_str
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(100, 100, 100)

# ─── MAIN HANDLER ─────────────────────────────────────────────────────────────

def handler(event: dict, context) -> None:  # noqa: ANN001
    """Entry point – invoked asynchronously. Each invocation processes one batch
    of slides, then chains to itself for the next batch via async Lambda invoke."""
    job_id: str = event["job_id"]
    structure: dict = event["structure"]
    settings: dict = event.get("settings", {})
    api_key: str = settings.get("api_key", "")
    logo_url: str = settings.get("logo_url", "")

    # Lambda-chaining state — absent on the first invocation
    batch_num: int = event.get("batch_num", 1)
    total_batches: int = event.get("total_batches", 0)   # computed on first call
    wip_s3_key: str | None = event.get("wip_s3_key")

    is_first = batch_num == 1
    if is_first:
        print(f"[{job_id}] Build-slides started. {len(structure.get('slides', []))} slides.")
    else:
        print(f"[{job_id}] Batch {batch_num}/{total_batches} invocation started.")

    client = OpenAI(api_key=api_key, base_url=QWEN_BASE_URL, timeout=600.0)
    s3 = boto3.client("s3")

    try:
        # Pre-compute which slides belong to this batch so downloads are scoped
        slides: list = structure.get("slides", [])
        batch_slides = slides[(batch_num - 1) * BATCH_SIZE : batch_num * BATCH_SIZE]

        # ── Always: download logo ──────────────────────────────────────────
        print(f"[{job_id}] Stage 0: Checking for logo...")
        logo_bytes: bytes | None = None
        if logo_url:
            try:
                req = urlrequest.Request(url=logo_url, method="GET")
                with urlrequest.urlopen(req, timeout=15) as resp:
                    logo_bytes = resp.read()
                print(f"[{job_id}] Stage 0: Logo downloaded ({len(logo_bytes)} bytes)")
            except Exception:  # noqa: BLE001
                logo_bytes = None

        # ── Always: pre-download background images ─────────────────────────
        print(f"[{job_id}] Stage 0.2: Pre-downloading background images...")
        image_urls: set[str] = set()
        for slide in batch_slides:
            url = slide.get("image_url")
            if url:
                image_urls.add(url)
        image_buffers: dict[str, bytes] = {}
        if image_urls:
            print(f"[{job_id}] Stage 0.2: Downloading {len(image_urls)} unique images...")
            for url in sorted(image_urls):
                try:
                    req = urlrequest.Request(url=url, method="GET")
                    with urlrequest.urlopen(req, timeout=15) as resp:
                        image_buffers[url] = resp.read()
                except Exception:  # noqa: BLE001
                    print(f"[{job_id}] Stage 0.2: Failed to download image {url[:80]}..., skipping")
            print(f"[{job_id}] Stage 0.2: Downloaded {len(image_buffers)}/{len(image_urls)} images")
        else:
            print(f"[{job_id}] Stage 0.2: No background images in structure")

        # ── Always: download icons ─────────────────────────────────────────
        print(f"[{job_id}] Stage 0.5: Collecting icons from structure...")
        icons: dict[str, bytes] = {}
        icon_names: set[str] = set()
        for slide in batch_slides:
            for col in slide.get("columns", []):
                for bullet in col.get("bullets", []):
                    icon_name = bullet.get("icon", "")
                    if icon_name:
                        icon_names.add(icon_name)
        if icon_names:
            primary_color = settings.get("primary_color", "#C00000")
            print(f"[{job_id}] Stage 0.5: Downloading {len(icon_names)} unique icons (PNG)...")
            for name in sorted(icon_names):
                try:
                    buf = io.BytesIO()
                    png_name = name.replace(".svg", ".png")
                    s3.download_fileobj(S3_OUTPUT_BUCKET, f"icons/{png_name}", buf)
                    # Recolour PNG to the user's primary colour
                    recolored = _recolor_png(buf.getvalue(), primary_color)
                    icons[name] = recolored
                except Exception:  # noqa: BLE001
                    print(f"[{job_id}] Stage 0.5: Failed to download icon {name}, skipping")
            print(f"[{job_id}] Stage 0.5: Downloaded {len(icons)}/{len(icon_names)} icons")
        else:
            print(f"[{job_id}] Stage 0.5: No icons referenced in structure")
        has_icons = len(icons) > 0

        # ── First invocation only: one-time setup ─────────────────────────
        if is_first:
            # Null out any image URLs that failed so the AI won't reference them
            for slide in slides:
                url = slide.get("image_url")
                if url and url not in image_buffers:
                    slide["image_url"] = None

            # Assign section numbers to divider slides
            section_num = 0
            for slide in slides:
                if slide.get("layout") == "section_divider":
                    section_num += 1
                    slide["section_number"] = section_num

            total_slides = len(slides)
            total_batches = math.ceil(total_slides / BATCH_SIZE)
            print(f"[{job_id}] {total_slides} slides → {total_batches} batches of {BATCH_SIZE}")
            _update_job(job_id, status="building", stage_message=f"Building slides 1-{min(BATCH_SIZE, total_slides)} of {total_slides}…")

        # ── Load or create prs ─────────────────────────────────────────────
        from pptx import Presentation as _Prs
        import pptx.util
        if wip_s3_key:
            print(f"[{job_id}] Loading WIP pptx from {wip_s3_key}...")
            wip_buf = io.BytesIO()
            s3.download_fileobj(S3_OUTPUT_BUCKET, wip_s3_key, wip_buf)
            wip_buf.seek(0)
            prs = _Prs(wip_buf)
        else:
            prs = _Prs()
            prs.slide_width = pptx.util.Inches(13.33)
            prs.slide_height = pptx.util.Inches(7.5)

        # ── Process this batch ─────────────────────────────────────────────
        total_slides = len(slides)
        fn_name = f"build_batch_{batch_num}"
        start_slide = (batch_num - 1) * BATCH_SIZE + 1
        end_slide = min(batch_num * BATCH_SIZE, total_slides)
        has_logo = logo_bytes is not None
        namespace = _make_namespace(image_buffers)

        print(f"[{job_id}] Batch {batch_num}/{total_batches}: slides {start_slide}-{end_slide}")
        _update_job(job_id, stage_message=f"Building slides {start_slide}-{end_slide} of {total_slides}…")

        batch_code = _build_batch_code(
            batch_slides, batch_num, total_batches, settings, client, job_id, has_logo, has_icons,
        )

        # Execute with up to 3 self-correction attempts
        batch_ok = False
        last_error = ""
        for attempt in range(3):
            slide_count_before = len(prs.slides._sldIdLst)
            try:
                _execute_batch(batch_code, prs, namespace, logo_bytes, fn_name, icons, batch_slides)
                print(f"[{job_id}] Batch {batch_num}/{total_batches}: executed on attempt {attempt + 1}")
                batch_ok = True
                break
            except Exception as exc:  # noqa: BLE001
                rolled_back = _rollback_slides(prs, slide_count_before)
                print(f"[{job_id}] Rolled back {rolled_back} orphan slide(s) from failed attempt")
                tb = traceback.format_exc()
                print(f"[{job_id}] Batch {batch_num}/{total_batches}: attempt {attempt + 1} FAILED")
                print(f"[{job_id}] Error: {exc}")
                print(f"[{job_id}] Traceback:\n{tb}")
                # Log first 500 chars of failing code for pattern analysis
                code_preview = batch_code[:500].replace("\n", "\\n")
                print(f"[{job_id}] Failing code preview: {code_preview}...")
                last_error = str(exc)
                if attempt < 2:
                    _update_job(job_id, stage_message=f"Fixing batch {batch_num} (attempt {attempt + 2}/3)…")
                    print(f"[{job_id}] Batch {batch_num}/{total_batches}: requesting fix from API...")
                    fix_code = batch_code
                    fix_exc = None
                    for fix_retry in range(3):
                        try:
                            print(f"[{job_id}] Batch {batch_num}/{total_batches}: fix API call (attempt {fix_retry+1}/3)...")
                            fix_response = client.chat.completions.create(
                                model=QWEN_MODEL,
                                messages=[
                                    {"role": "system", "content": _BUILD_SYSTEM},
                                    {"role": "user", "content": (
                                        f"The following code raised an error:\n\n{batch_code}\n\n"
                                        f"Error: {last_error}\n\n"
                                        f"Fix the code and return only the corrected function named `{fn_name}`.\n"
                                        "REMINDER: All modules are already injected as global variables \u2014 "
                                        "remove any import statements and use the pre-injected names directly."
                                    )},
                                ],
                                max_tokens=16000,
                                timeout=600.0,
                            )
                            fix_code = fix_response.choices[0].message.content or batch_code
                            print(f"[{job_id}] Batch {batch_num}/{total_batches}: fix received ({len(fix_code)} chars)")
                            fix_exc = None
                            break
                        except Exception as fix_err:  # noqa: BLE001
                            fix_exc = fix_err
                            print(f"[{job_id}] Batch {batch_num}/{total_batches}: fix API attempt {fix_retry+1} failed: {fix_err}")
                            if fix_retry < 2:
                                wait = (fix_retry + 1) * 10
                                print(f"[{job_id}] Waiting {wait}s before fix retry...")
                                time.sleep(wait)
                    if fix_exc is not None:
                        raise RuntimeError(
                            f"Fix API call failed after 3 attempts: {fix_exc}"
                        ) from fix_exc
                    batch_code = fix_code

        if not batch_ok:
            raise RuntimeError(
                f"Batch {batch_num}/{total_batches} failed after 3 attempts. Last error: {last_error}"
            )

        # ── Save batch code to S3 ──────────────────────────────────────────
        s3.put_object(
            Bucket=S3_OUTPUT_BUCKET,
            Key=f"wip/{job_id}/code_{batch_num}.txt",
            Body=batch_code.encode(),
        )
        # ── Add references slide on final batch ────────────────────────────────────
        if batch_num >= total_batches:
            print(f"[{job_id}] Adding references slide...")
            _add_references_slide(prs, structure, settings)
        # ── Serialize current prs to bytes ─────────────────────────────────
        prs_buf = io.BytesIO()
        prs.save(prs_buf)
        prs_bytes = prs_buf.getvalue()

        if batch_num < total_batches:
            # ── More batches: save WIP and chain to next invocation ────────
            new_wip_key = f"wip/{job_id}/pptx.pptx"
            s3.put_object(Bucket=S3_OUTPUT_BUCKET, Key=new_wip_key, Body=prs_bytes)
            print(f"[{job_id}] WIP saved ({len(prs_bytes)} bytes). Invoking batch {batch_num + 1}/{total_batches}...")
            boto3.client("lambda").invoke(
                FunctionName=context.function_name,
                InvocationType="Event",
                Payload=json.dumps({
                    "job_id": job_id,
                    "structure": structure,
                    "settings": settings,
                    "batch_num": batch_num + 1,
                    "total_batches": total_batches,
                    "wip_s3_key": new_wip_key,
                }).encode(),
            )
            print(f"[{job_id}] Invoked self for batch {batch_num + 1}/{total_batches}")

        else:
            # ── Last batch: collect code, upload final PPTX, mark done ────
            all_code: list[str] = []
            for i in range(1, total_batches + 1):
                try:
                    resp = s3.get_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/code_{i}.txt")
                    all_code.append(resp["Body"].read().decode())
                except Exception:  # noqa: BLE001
                    pass
            _update_job(job_id, pptx_code="\n\n# --- BATCH ---\n\n".join(all_code))
            print(f"[{job_id}] Final batch done. Uploading PPTX ({len(prs_bytes)} bytes)...")
            key = f"presentations/{job_id}.pptx"
            s3.put_object(
                Bucket=S3_OUTPUT_BUCKET,
                Key=key,
                Body=prs_bytes,
                ContentType=(
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.presentation"
                ),
            )
            download_url = s3.generate_presigned_url(
                "get_object",
                Params={"Bucket": S3_OUTPUT_BUCKET, "Key": key},
                ExpiresIn=3600,
            )
            _update_job(
                job_id,
                status="done",
                stage_message="Your presentation is ready!",
                download_url=download_url,
            )
            print(f"[{job_id}] DONE. Download: {download_url}")

            # Clean up WIP files
            for i in range(1, total_batches + 1):
                try:
                    s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/code_{i}.txt")
                except Exception:  # noqa: BLE001
                    pass
            try:
                s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/pptx.pptx")
            except Exception:  # noqa: BLE001
                pass

    except Exception as exc:  # noqa: BLE001
        print(f"[{job_id}] ERROR: {exc}")
        _update_job(
            job_id,
            status="error",
            stage_message="Something went wrong.",
            error_message=str(exc),
        )
        raise
