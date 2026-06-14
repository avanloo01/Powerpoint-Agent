"""
Agent-loop Lambda: orchestrates Research → Structure → Build for a PowerPoint presentation.
Invoked asynchronously by start_job. Updates a Supabase `jobs` row at each stage.
"""
from __future__ import annotations

import builtins as _builtins
import io
import json
import math
import os
import re
import textwrap
# import uuid
# from urllib import error as urlerror
from urllib import request as urlrequest

import boto3
from openai import OpenAI

# ─── CONSTANTS ────────────────────────────────────────────────────────────────

S3_OUTPUT_BUCKET = os.environ.get("OUTPUT_BUCKET", "")
QWEN_BASE_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1"
QWEN_MODEL = os.environ.get("QWEN_MODEL", "qwen3.6-plus")
SUPABASE_URL = os.environ.get("SUPABASE_URL", "").rstrip("/")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")

# ─── SUPABASE HELPERS ─────────────────────────────────────────────────────────

def _supabase_request(method: str, url: str, headers: dict, body: dict | None = None) -> object:
    payload = None
    request_headers = dict(headers)
    if body is not None:
        payload = json.dumps(body).encode("utf-8")
        request_headers["Content-Type"] = "application/json"
    req = urlrequest.Request(url=url, method=method, data=payload, headers=request_headers)
    try:
        with urlrequest.urlopen(req, timeout=10) as resp:
            raw = resp.read().decode("utf-8")
            return json.loads(raw) if raw else {}
    except Exception:  # noqa: BLE001 – best-effort status updates must never crash the loop
        return {}


def _update_job(job_id: str, **fields: object) -> None:
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        return
    url = f"{SUPABASE_URL}/rest/v1/jobs?id=eq.{job_id}"
    _supabase_request(
        "PATCH",
        url,
        {
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
            "Prefer": "return=minimal",
        },
        fields,  # type: ignore[arg-type]
    )


# ─── RESEARCH AGENT ──────────────────────────────────────────────────────────

def _research(prompt: str, client: OpenAI, job_id: str) -> str:
    """Stage 1 – Use Qwen with web search to gather current facts and data."""
    _update_job(job_id, status="researching", stage_message="Researching your topic\u2026")

    system = (
        "You are a research analyst preparing materials for a business presentation. "
        "Use web search to gather current facts, statistics, trends and examples. "
        "Produce a thorough markdown document covering: key facts and data, recent trends, "
        "notable examples or case studies, and key takeaways. "
        "Include specific numbers and dates where available."
    )
    response = client.chat.completions.create(
        model=QWEN_MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": f"Research this topic for a business presentation:\n\n{prompt}"},
        ],
        extra_body={"enable_search": True},
        max_tokens=4000,
    )
    return response.choices[0].message.content or ""


# ─── STRUCTURE AGENT ─────────────────────────────────────────────────────────

_STRUCTURE_SCHEMA = textwrap.dedent("""\
Return ONLY a valid JSON object matching this schema (no markdown fences):
{
  "presentation_title": "string",
  "slides": [
    {
      "slide_title": "string",
      "section_label": "string",
      "layout": "two_columns | three_columns | full_width",
      "columns": [
        {
          "width_ratio": <0.33 | 0.5 | 0.67 | 1.0>,
          "box_header": "string or null",
          "subtitle": "string or null",
          "content_type": "chart | bullet_list | text | news_cards | icon_grid",
          "chart": {
            "chart_type": "bar | grouped_bar | line | pie",
            "title": "string",
            "x_labels": ["string"],
            "series": [{"name": "string", "values": [<number>]}]
          },
          "bullets": [
            {"icon": "acorn.svg | address-book.svg | air-traffic-control.svg | airplane-landing.svg | airplane-takeoff.svg | airplane-tilt.svg | airplay.svg | alarm.svg | alien.svg | align-bottom.svg | align-center-horizontal.svg | align-center-vertical.svg | align-left.svg | align-right.svg | align-top.svg | amazon-logo.svg | ambulance.svg | anchor-simple.svg | android-logo.svg | angle.svg | aperture.svg | app-store-logo.svg | apple-logo.svg | approximate-equals.svg | archive.svg | armchair.svg | arrow-arc-left.svg | arrow-arc-right.svg | arrow-bend-down-left.svg | arrow-bend-down-right.svg | arrow-bend-left-down.svg | arrow-bend-left-up.svg | arrow-bend-right-down.svg | arrow-bend-right-up.svg | arrow-bend-up-left.svg | arrow-bend-up-right.svg | arrow-clockwise.svg | arrow-counter-clockwise.svg | arrow-down-left.svg | arrow-down-right.svg | arrow-down.svg | arrow-fat-down.svg | arrow-fat-left.svg | arrow-fat-right.svg | arrow-fat-up.svg | arrow-left.svg | arrow-up-left.svg | arrow-up-right.svg | arrow-up.svg | arrows-clockwise.svg | arrows-down-up.svg | arrows-horizontal.svg | arrows-in-line-horizontal.svg | arrows-in-line-vertical.svg | arrows-in.svg | arrows-left-right.svg | arrows-merge.svg | arrows-out-cardinal.svg | arrows-out-line-horizontal.svg | arrows-out-line-vertical.svg | arrows-out.svg | arrows-split.svg | arrows-vertical.svg | asterisk.svg | at.svg | atom.svg | axe.svg | baby-carriage.svg | baby.svg | backspace.svg | bag.svg | balloon.svg | bandaids.svg | bank.svg | barbell.svg | barcode.svg | barn.svg | barricade.svg | baseball-helmet.svg | basket.svg | basketball.svg | bathtub.svg | battery-charging.svg | battery-high.svg | battery-warning.svg | beanie.svg | bed.svg | beer-bottle.svg | beer-stein.svg | bell-ringing.svg | bell-slash.svg | bell.svg | belt.svg | bezier-curve.svg | bicycle.svg | binary.svg | binoculars.svg | biohazard.svg | bird.svg | blueprint.svg | bluetooth-slash.svg | bluetooth.svg | boat.svg | bomb.svg | bone.svg | book-bookmark.svg | book.svg | books.svg | bounding-box.svg | bowl-steam.svg | bowling-ball.svg | box-arrow-down.svg | box-arrow-up.svg | boxing-glove.svg | brackets-angle.svg | brackets-curly.svg | brackets-round.svg | brackets-square.svg | brain.svg | brandy.svg | bread.svg | bridge.svg | broadcast.svg | broom.svg | bug.svg | buildings.svg | bulldozer.svg | butterfly.svg | cable-car.svg | cactus.svg | cake.svg | calculator.svg | calendar-blank.svg | calendar-check.svg | calendar-minus.svg | calendar-plus.svg | calendar-slash.svg | call-bell.svg | camera-plus.svg | camera-rotate.svg | camera-slash.svg | camera.svg | car-profile.svg | carrot.svg | cash-register.svg | cassette-tape.svg | castle-turret.svg | cat.svg | cell-signal-full.svg | cell-signal-slash.svg | cell-tower.svg | certificate.svg | chalkboard-simple.svg | champagne.svg | charging-station.svg | chart-bar.svg | chart-line-down.svg | chart-line-up.svg | chart-line.svg | chart-pie.svg | chart-scatter.svg | chat-centered-dots.svg | chat-centered-slash.svg | chat-text.svg | chats.svg | cheers.svg | cheese.svg | chef-hat.svg | cherries.svg | church.svg | cigarette-slash.svg | cigarette.svg | circle.svg | circuitry.svg | city.svg | clipboard-text.svg | clock-clockwise.svg | clock-counter-clockwise.svg | clock.svg | cloud-arrow-down.svg | cloud-arrow-up.svg | cloud-check.svg | cloud-lightning.svg | cloud-moon.svg | cloud-rain.svg | cloud-slash.svg | cloud-snow.svg | cloud-sun.svg | cloud.svg | clover.svg | club.svg | coat-hanger.svg | code-simple.svg | code.svg | coffee.svg | coins.svg | columns-plus-left.svg | columns-plus-right.svg | columns.svg | compass-rose.svg | compass-tool.svg | compass.svg | confetti.svg | contactless-payment.svg | cookie.svg | cooking-pot.svg | copy.svg | copyright.svg | corners-in.svg | corners-out.svg | couch.svg | cow.svg | cowboy-hat.svg | cpu.svg | crane-tower.svg | cricket.svg | crop.svg | cross.svg | crosshair-simple.svg | crown-simple.svg | cube-focus.svg | cube.svg | currency-btc.svg | currency-cny.svg | currency-dollar-simple.svg | currency-dollar.svg | currency-eth.svg | currency-eur.svg | currency-gbp.svg | currency-inr.svg | currency-jpy.svg | currency-krw.svg | currency-kzt.svg | currency-ngn.svg | currency-rub.svg | cursor-text.svg | cursor.svg | database.svg | desktop-tower.svg | detective.svg | device-mobile-slash.svg | device-mobile.svg | device-rotate.svg | diamond.svg | diamonds-four.svg | dice-six.svg | disco-ball.svg | discord-logo.svg | divide.svg | dna.svg | dog.svg | door.svg | download-simple.svg | dress.svg | dresser.svg | drone.svg | drop-simple.svg | drop-slash.svg | dropbox-logo.svg | egg-crack.svg | eject.svg | envelope-simple-open.svg | envelope-simple.svg | equalizer.svg | equals.svg | eraser.svg | escalator-down.svg | escalator-up.svg | exclamation-mark.svg | export.svg | eye-slash.svg | eye.svg | eyedropper.svg | face-mask.svg | facebook-logo.svg | factory.svg | faders.svg | fan.svg | fast-forward.svg | feather.svg | figma-logo.svg | file-archive.svg | file-arrow-down.svg | file-arrow-up.svg | file-cloud.svg | file-code.svg | file-pdf.svg | file-text.svg | film-slate.svg | fingerprint.svg | fire-extinguisher.svg | fire-simple.svg | first-aid-kit.svg | fish-simple.svg | flag.svg | flashlight.svg | flask.svg | flower-tulip.svg | flying-saucer.svg | folder-simple.svg | football.svg | fork-knife.svg | function.svg | funnel.svg | game-controller.svg | gas-pump.svg | gavel.svg | gear-six.svg | gender-female.svg | gender-intersex.svg | gender-male.svg | gender-neuter.svg | gender-nonbinary.svg | gender-transgender.svg | ghost.svg | gift.svg | git-branch.svg | git-commit.svg | git-diff.svg | git-fork.svg | git-merge.svg | git-pull-request.svg | github-logo.svg | gitlab-logo-simple.svg | globe-hemisphere-east.svg | globe-hemisphere-west.svg | globe.svg | google-chrome-logo.svg | google-drive-logo.svg | google-logo.svg | google-photos-logo.svg | google-play-logo.svg | graduation-cap.svg | graph.svg | graphics-card.svg | greater-than-or-equal.svg | greater-than.svg | guitar.svg | hair-dryer.svg | hamburger.svg | hammer.svg | hand-eye.svg | hand-fist.svg | hand-heart.svg | hand-peace.svg | hand-pointing.svg | hand-soap.svg | hand-waving.svg | hand.svg | handbag-simple.svg | hands-clapping.svg | hands-praying.svg | handshake.svg | hard-hat.svg | hash.svg | head-circuit.svg | headlights.svg | headphones.svg | heart-break.svg | heart.svg | heartbeat.svg | hexagon.svg | high-heel.svg | highlighter.svg | hockey.svg | horse.svg | hospital.svg | hourglass-medium.svg | house.svg | ice-cream.svg | identification-card.svg | image.svg | infinity.svg | info.svg | instagram-logo.svg | intersect-three.svg | island.svg | jar-label.svg | joystick.svg | kanban.svg | key.svg | keyhole.svg | knife.svg | ladder-simple.svg | lamp.svg | layout.svg | leaf.svg | less-than-or-equal.svg | less-than.svg | lightbulb.svg | lightning-slash.svg | lightning.svg | line-segments.svg | link-break.svg | link.svg | linkedin-logo.svg | linux-logo.svg | list-bullets.svg | list-checks.svg | lock-simple-open.svg | lock-simple.svg | magnet-straight.svg | magnifying-glass.svg | mailbox.svg | map-pin.svg | martini.svg | mask-happy.svg | mask-sad.svg | math-operations.svg | medal-military.svg | megaphone-simple.svg | memory.svg | messenger-logo.svg | meta-logo.svg | metronome.svg | microphone-slash.svg | microphone.svg | microscope.svg | minus.svg | money.svg | moon.svg | moped.svg | mosque.svg | mountains.svg | mouse.svg | music-notes-simple.svg | navigation-arrow.svg | newspaper.svg | not-equals.svg | note-blank.svg | nuclear-plant.svg | number-eight.svg | number-five.svg | number-four.svg | number-nine.svg | number-one.svg | number-seven.svg | number-six.svg | number-three.svg | number-two.svg | number-zero.svg | numpad.svg | octagon.svg | office-chair.svg | open-ai-logo.svg | oven.svg | package.svg | palette.svg | paper-plane-right.svg | paperclip.svg | parallelogram.svg | park.svg | pause.svg | paw-print.svg | paypal-logo.svg | peace.svg | pen-nib.svg | pencil-simple-slash.svg | pencil-simple.svg | pentagon.svg | pentagram.svg | pepper.svg | percent.svg | person-simple-bike.svg | person-simple-hike.svg | person-simple-run.svg | person-simple-ski.svg | person-simple-snowboard.svg | person-simple-swim.svg | person-simple-tai-chi.svg | person-simple-throw.svg | person.svg | phone-call.svg | phone-disconnect.svg | phone-incoming.svg | phone-outgoing.svg | phone-slash.svg | phone.svg | pi.svg | piano-keys.svg | picnic-table.svg | picture-in-picture.svg | piggy-bank.svg | pill.svg | ping-pong.svg | pint-glass.svg | pipe-wrench.svg | pipe.svg | pizza.svg | placeholder.svg | planet.svg | plant.svg | play.svg | playlist.svg | plug.svg | plugs.svg | plus-minus.svg | plus.svg | poker-chip.svg | polygon.svg | popcorn.svg | popsicle.svg | potted-plant.svg | power.svg | prescription.svg | presentation-chart.svg | printer.svg | prohibit.svg | pulse.svg | push-pin-slash.svg | push-pin.svg | puzzle-piece.svg | qr-code.svg | question-mark.svg | quotes.svg | rabbit.svg | racquet.svg | radical.svg | radio.svg | radioactive.svg | rainbow-cloud.svg | ranking.svg | receipt.svg | record.svg | rectangle-dashed.svg | rectangle.svg | recycle.svg | reddit-logo.svg | repeat-once.svg | repeat.svg | resize.svg | rewind.svg | road-horizon.svg | robot.svg | rocket-launch.svg | rows-plus-bottom.svg | rows-plus-top.svg | rows.svg | rss.svg | ruler.svg | sailboat.svg | scales.svg | scan-smiley.svg | scissors.svg | scooter.svg | screencast.svg | screwdriver.svg | seal-check.svg | seat.svg | security-camera.svg | share-network.svg | shield-check.svg | shield-plus.svg | shield-slash.svg | shield-star.svg | shield-warning.svg | shield.svg | shirt-folded.svg | shopping-cart-simple.svg | shovel.svg | shower.svg | shrimp.svg | shuffle-simple.svg | sigma.svg | sign-in.svg | sign-out.svg | signature.svg | siren.svg | skip-back.svg | skip-forward.svg | skull.svg | skype-logo.svg | slack-logo.svg | sliders.svg | smiley-angry.svg | smiley-meh.svg | smiley-melting.svg | smiley-nervous.svg | smiley-sad.svg | smiley-wink.svg | smiley-x-eyes.svg | smiley.svg | snapchat-logo.svg | sneaker.svg | snowflake.svg | soccer-ball.svg | sock.svg | solar-panel.svg | soundcloud-logo.svg | spade.svg | sparkle.svg | speaker-hifi.svg | speaker-simple-high.svg | speaker-simple-none.svg | speaker-simple-slash.svg | speaker-simple-x.svg | speedometer.svg | spinner-gap.svg | spiral.svg | spotify-logo.svg | spray-bottle.svg | square-logo.svg | square.svg | stack-overflow-logo.svg | stack.svg | star-and-crescent.svg | star-of-david.svg | star.svg | steam-logo.svg | steps.svg | stethoscope.svg | stool.svg | storefront.svg | strategy.svg | stripe-logo.svg | student.svg | suitcase-rolling.svg | sun.svg | sunglasses.svg | swimming-pool.svg | sword.svg | synagogue.svg | syringe.svg | t-shirt.svg | table.svg | tag.svg | taxi.svg | telegram-logo.svg | television-simple.svg | tennis-ball.svg | tent.svg | terminal.svg | test-tube.svg | textbox.svg | thermometer-simple.svg | threads-logo.svg | thumbs-down.svg | thumbs-up.svg | ticket.svg | tiktok-logo.svg | tilde.svg | tip-jar.svg | tipi.svg | toggle-right.svg | toilet-paper.svg | toilet.svg | tooth.svg | tractor.svg | trademark-registered.svg | trademark.svg | traffic-cone.svg | traffic-signal.svg | train-simple.svg | translate.svg | trash.svg | tree-evergreen.svg | tree-palm.svg | tree-structure.svg | tree-view.svg | tree.svg | trend-down.svg | trend-up.svg | triangle-dashed.svg | triangle.svg | trophy.svg | truck.svg | tumblr-logo.svg | twitch-logo.svg | twitter-logo.svg | umbrella-simple.svg | upload-simple.svg | upload.svg | usb.svg | user-check.svg | user-focus.svg | user-list.svg | user-minus.svg | user-plus.svg | user.svg | users-three.svg | vault.svg | vector-three.svg | vector-two.svg | vibrate.svg | video-camera-slash.svg | video-camera.svg | video.svg | vignette.svg | vinyl-record.svg | virus.svg | voicemail.svg | volleyball.svg | wall.svg | wallet.svg | warning.svg | washing-machine.svg | watch.svg | wave-sawtooth.svg | wave-sine.svg | wave-square.svg | waves.svg | webcam-slash.svg | webcam.svg | wechat-logo.svg | whatsapp-logo.svg | wheelchair.svg | wifi-high.svg | wifi-slash.svg | wind.svg | windmill.svg | windows-logo.svg | wine.svg | wrench.svg | x-logo.svg | x.svg | yarn.svg | yin-yang.svg | youtube-logo.svg",
             "title": "string", "description": "string"}
          ],
          "news_cards": [
            {"date": "string", "source": "string", "headline": "string", "subtext": "string"}
          ],
          "text": "string"
        }
      ],
      "column_separator": "line | causal_line | none",
      "conclusion_box": "string or null",
      "sources": "string or null"
    }
  ]
}
Rules:
- two_columns: exactly 2 columns, width_ratios must sum to 1.0
- three_columns: exactly 3 columns, each width_ratio = 0.33
- full_width: exactly 1 column, width_ratio = 1.0
- For charts: always include realistic data matching the research findings
- Use 4–8 slides total; group related slides under the same section_label
""")


def _structure(prompt: str, research_md: str, client: OpenAI, job_id: str) -> dict:
    """Stage 2 – Design the presentation structure as a validated JSON blueprint."""
    _update_job(job_id, status="structuring", stage_message="Building presentation structure\u2026")

    system = (
        "You are a McKinsey-trained presentation designer. "
        "Based on the research findings, design a compelling, insight-driven PowerPoint presentation. "
        + _STRUCTURE_SCHEMA
    )
    response = client.chat.completions.create(
        model=QWEN_MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": f"Topic: {prompt}\n\nResearch:\n{research_md}"},
        ],
        response_format={"type": "json_object"},
        max_tokens=6000,
    )
    raw = response.choices[0].message.content or "{}"
    return json.loads(raw)


# ─── BUILD AGENT ─────────────────────────────────────────────────────────────

_BUILD_SYSTEM = textwrap.dedent("""\
You are an expert python-pptx developer. Write a Python function called
`build_presentation(prs, logo_bytes=None)` that adds all slides to the given
python-pptx Presentation object `prs`.

CONSTRAINTS:
- `prs` already has slide_width = Inches(13.33), slide_height = Inches(7.5)
- Use `prs.slide_layouts[6]` (blank layout) for every slide
- DO NOT call Presentation() yourself — use the `prs` argument
- Available names in scope: Inches, Pt, Emu, Cm, RGBColor, PP_ALIGN,
  ChartData, XL_CHART_TYPE, io, math, json
- Return ONLY the Python function code, no markdown fences, no other text
- The function must be syntactically valid Python 3.12

STYLE GUIDE:
- Slide background: white
- Section label: top-left, 9 pt, RGB(128,128,128)
- Slide title: bold, 22 pt, black, positioned below section label
- Box headers: filled rectangles with white bold text (12 pt); use primary color; Except for a box header that fills one column, the headers should not always take up the entire width of the column. This leaves appropriate padding. For example, for two columns, the first is slightly shorter than the first column. and the second does not immediately start at the beginning of the second column, but after a small gap. The same applies to three columns or columns where we have two-thirds and one-third (or vice versa).
- Column separators: thin vertical line (0.5 pt, light gray) centered between columns
- Causal separator: same vertical line + small filled triangle pointing right
- Bullet lists: icon placeholder (filled circle) + bold title + description text
- Charts: use ChartData + slide.shapes.add_chart() for bar/line/pie charts
- Conclusion box: thin-bordered rectangle (1 pt, primary color), centered italic text 11 pt
- Sources: bottom-left, 8 pt, RGB(128,128,128)
- Logo: the function receives a keyword argument `logo_bytes` (bytes or None).
  When logo_bytes is not None, place the logo on EVERY slide at the top-right
  corner. Use io.BytesIO(logo_bytes) + slide.shapes.add_picture(). The logo
  must be small (0.5–0.7 inches tall, proportional width). Its right edge must
  align with the right edge of the rightmost column. Its top edge must align
  with the section label / slide title top edge.
""")


def _build_code(structure: dict, settings: dict, client: OpenAI, job_id: str, has_logo: bool = False) -> str:
    """Stage 3a – Generate python-pptx code from the structure blueprint."""
    _update_job(job_id, status="building", stage_message="Building your presentation\u2026")

    primary = settings.get("primary_color", "#C00000")
    accent = settings.get("accent_color", "#A6CAEC")
    logo_note = "A logo_bytes parameter WILL be provided — place the logo on every slide." if has_logo else "No logo will be provided."
    user_msg = (
        f"Primary color: {primary}\n"
        f"Accent color: {accent}\n"
        f"Logo: {logo_note}\n\n"
        f"Presentation structure:\n{json.dumps(structure, indent=2)}"
    )
    response = client.chat.completions.create(
        model=QWEN_MODEL,
        messages=[
            {"role": "system", "content": _BUILD_SYSTEM},
            {"role": "user", "content": user_msg},
        ],
        max_tokens=8000,
    )
    return response.choices[0].message.content or ""


# ─── EXECUTE GENERATED CODE ───────────────────────────────────────────────────

def _execute(code: str, logo_bytes: bytes | None = None) -> bytes:
    """
    Execute the AI-generated build_presentation(prs, logo_bytes=...) function
    in a restricted namespace and return the resulting PPTX bytes.

    Security notes:
    - __builtins__ is replaced with a curated allowlist; __import__, open,
      eval, and exec are excluded.
    - os, sys, socket are not injected, preventing system/network access.
    - python-pptx itself runs in its own module scope with full builtins, which
      is fine because it is a vetted dependency, not user-supplied code.
    """
    import pptx
    import pptx.chart.data
    import pptx.dml.color
    import pptx.enum.chart
    import pptx.enum.text
    import pptx.util
    from pptx import Presentation as _Prs

    # Strip markdown fences the model may have added
    code = re.sub(r"```(?:python)?\s*", "", code).strip().rstrip("`").strip()

    _SAFE_NAMES = (
        "abs", "bool", "dict", "enumerate", "float", "hasattr", "int",
        "isinstance", "len", "list", "max", "min", "print", "range",
        "round", "set", "str", "sum", "tuple", "zip",
    )
    safe_builtins = {
        name: getattr(_builtins, name)
        for name in _SAFE_NAMES
        if hasattr(_builtins, name)
    }

    namespace: dict = {
        "__builtins__": safe_builtins,
        # pptx utilities
        "Inches": pptx.util.Inches,
        "Pt": pptx.util.Pt,
        "Emu": pptx.util.Emu,
        "Cm": pptx.util.Cm,
        "RGBColor": pptx.dml.color.RGBColor,
        "PP_ALIGN": pptx.enum.text.PP_ALIGN,
        "ChartData": pptx.chart.data.ChartData,
        "XL_CHART_TYPE": pptx.enum.chart.XL_CHART_TYPE,
        # safe stdlib modules
        "io": io,
        "math": math,
        "json": json,
    }

    exec(code, namespace)  # noqa: S102

    build_fn = namespace.get("build_presentation")
    if not callable(build_fn):
        raise ValueError("Generated code does not define a callable 'build_presentation'")

    prs = _Prs()
    prs.slide_width = pptx.util.Inches(13.33)
    prs.slide_height = pptx.util.Inches(7.5)
    build_fn(prs, logo_bytes=logo_bytes)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─── MAIN HANDLER ─────────────────────────────────────────────────────────────

def handler(event: dict, context) -> None:  # noqa: ANN001
    """Entry point – invoked asynchronously; no HTTP response required."""
    job_id: str = event["job_id"]
    prompt: str = event["prompt"]
    settings: dict = event.get("settings", {})
    api_key: str = settings.get("api_key", "")

    client = OpenAI(api_key=api_key, base_url=QWEN_BASE_URL)
    s3 = boto3.client("s3")

    try:
        # ── Stage 0: Download logo (if provided) ────────────────────────────
        logo_bytes: bytes | None = None
        logo_url = settings.get("logo_url", "")
        if logo_url:
            try:
                req = urlrequest.Request(url=logo_url, method="GET")
                with urlrequest.urlopen(req, timeout=15) as resp:
                    logo_bytes = resp.read()
            except Exception:  # noqa: BLE001 – non-critical; proceed without logo
                logo_bytes = None

        # ── Stage 1: Research ───────────────────────────────────────────────
        research_md = _research(prompt, client, job_id)
        _update_job(job_id, research_md=research_md)

        # ── Stage 2: Structure ──────────────────────────────────────────────
        structure = _structure(prompt, research_md, client, job_id)
        _update_job(job_id, structure_md=json.dumps(structure, indent=2))

        # ── Stage 3a: Generate code ─────────────────────────────────────────
        pptx_code = _build_code(structure, settings, client, job_id, has_logo=logo_bytes is not None)
        _update_job(job_id, pptx_code=pptx_code)

        # ── Stage 3b: Execute with up to 3 self-correction attempts ─────────
        pptx_bytes: bytes | None = None
        last_error = ""

        for attempt in range(3):
            try:
                pptx_bytes = _execute(pptx_code, logo_bytes=logo_bytes)
                break
            except Exception as exc:  # noqa: BLE001
                last_error = str(exc)
                if attempt < 2:
                    fix_response = client.chat.completions.create(
                        model=QWEN_MODEL,
                        messages=[
                            {"role": "system", "content": _BUILD_SYSTEM},
                            {"role": "user", "content": (
                                f"The following code raised an error:\n\n{pptx_code}\n\n"
                                f"Error: {last_error}\n\n"
                                "Fix the code and return only the corrected function."
                            )},
                        ],
                        max_tokens=8000,
                    )
                    pptx_code = fix_response.choices[0].message.content or pptx_code
                    _update_job(job_id, pptx_code=pptx_code)

        if pptx_bytes is None:
            raise RuntimeError(f"Code execution failed after 3 attempts. Last error: {last_error}")

        # ── Upload to S3 ────────────────────────────────────────────────────
        key = f"presentations/{job_id}.pptx"
        s3.put_object(
            Bucket=S3_OUTPUT_BUCKET,
            Key=key,
            Body=pptx_bytes,
            ContentType=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )
        download_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": S3_OUTPUT_BUCKET, "Key": key},
            ExpiresIn=3600,
        )

        _update_job(
            job_id,
            status="done",
            stage_message="Your presentation is ready!",
            download_url=download_url,
        )

    except Exception as exc:  # noqa: BLE001
        _update_job(
            job_id,
            status="error",
            stage_message="Something went wrong.",
            error_message=str(exc),
        )
        raise
