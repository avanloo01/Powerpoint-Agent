"""
Lambda handler for generating a PowerPoint presentation using Qwen AI and python-pptx.
"""
from __future__ import annotations

import io
import json
import os
import re
import uuid
from urllib import error as urlerror
from urllib import parse as urlparse
from urllib import request as urlrequest

import boto3
from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor

S3_OUTPUT_BUCKET = os.environ.get("OUTPUT_BUCKET", "")
QWEN_BASE_URL = "https://dashscope-intl.aliyuncs.com/compatible-mode/v1"
QWEN_MODEL = os.environ.get("QWEN_MODEL", "qwen3.6-plus")
SUPABASE_URL = os.environ.get("SUPABASE_URL", "").rstrip("/")
SUPABASE_ANON_KEY = os.environ.get("SUPABASE_ANON_KEY", "")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")
SUPABASE_SETTINGS_TABLE = os.environ.get("SUPABASE_SETTINGS_TABLE", "user_settings")


def _extract_bearer_token(event: dict) -> str:
    headers = event.get("headers") or {}
    auth_header = headers.get("authorization") or headers.get("Authorization") or ""
    if not auth_header.startswith("Bearer "):
        return ""
    return auth_header.split(" ", 1)[1].strip()


def _supabase_request(method: str, url: str, headers: dict, body: dict | None = None) -> dict:
    payload = None
    request_headers = dict(headers)
    if body is not None:
        payload = json.dumps(body).encode("utf-8")
        request_headers["Content-Type"] = "application/json"

    req = urlrequest.Request(url=url, method=method, data=payload, headers=request_headers)
    try:
        with urlrequest.urlopen(req) as resp:
            raw = resp.read().decode("utf-8")
            return json.loads(raw) if raw else {}
    except urlerror.HTTPError as exc:
        detail = exc.read().decode("utf-8")
        raise RuntimeError(f"Supabase request failed ({exc.code}): {detail}") from exc


def _get_authenticated_user(token: str) -> dict:
    if not token:
        raise ValueError("Missing bearer token")
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        raise RuntimeError("SUPABASE_URL and SUPABASE_ANON_KEY must be configured")

    url = f"{SUPABASE_URL}/auth/v1/user"
    return _supabase_request(
        "GET",
        url,
        {
            "apikey": SUPABASE_ANON_KEY,
            "Authorization": f"Bearer {token}",
        },
    )


def _get_user_settings(user_id: str) -> dict | None:
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        raise RuntimeError("SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY must be configured")

    query = urlparse.urlencode(
        {
            "user_id": f"eq.{user_id}",
            "select": "api_key,primary_color,accent_color,logo_url",
            "limit": "1",
        }
    )
    url = f"{SUPABASE_URL}/rest/v1/{SUPABASE_SETTINGS_TABLE}?{query}"
    data = _supabase_request(
        "GET",
        url,
        {
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        },
    )

    if isinstance(data, list) and data:
        return data[0]
    return None


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a CSS hex color (e.g. '#4f46e5') to an RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    r, g, b = (int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def _build_presentation(
    slides_data: list[dict],
    primary_color: str,
    accent_color: str,
    logo_url: str | None,
) -> bytes:
    """Build a PPTX in memory and return the raw bytes."""
    prs = Presentation()

    primary_rgb = _hex_to_rgb(primary_color)
    accent_rgb = _hex_to_rgb(accent_color)

    for i, slide_data in enumerate(slides_data):
        # Use title+content layout for body slides, title-only for the first
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "")
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = primary_rgb
                    run.font.bold = True

        # Body / content
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            content = slide_data.get("content", "")
            lines = content.split("\n") if content else []
            for j, line in enumerate(lines):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = line.strip("- ").strip()
                for run in para.runs:
                    run.font.color.rgb = accent_rgb if j % 2 == 1 else RGBColor(0x11, 0x18, 0x27)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _parse_slides(raw: str) -> list[dict]:
    """
    Parse the AI response into a list of slide dicts.
    Accepts either a JSON array directly or JSON embedded in a markdown code block.
    """
    # Strip markdown fences if present
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).strip().rstrip("`").strip()
    try:
        data = json.loads(cleaned)
        if isinstance(data, list):
            return data
        if isinstance(data, dict) and "slides" in data:
            return data["slides"]
    except json.JSONDecodeError:
        pass

    # Fallback: create a single slide with the raw text
    return [{"title": "Presentation", "content": raw}]


def handler(event: dict, context) -> dict:  # noqa: ANN001
    """AWS Lambda entry point."""
    try:
        token = _extract_bearer_token(event)
        user = _get_authenticated_user(token)
        user_id = (user or {}).get("id", "")
        if not user_id:
            return _response(401, {"error": "Unauthorized"})

        settings = _get_user_settings(user_id) or {}

        body = json.loads(event.get("body") or "{}")
        prompt: str = body.get("prompt", "").strip()
        file_ids: list[str] = body.get("fileIDs") or []
        api_key: str = (settings.get("api_key") or "").strip()
        primary_color: str = settings.get("primary_color") or "#4f46e5"
        accent_color: str = settings.get("accent_color") or "#f59e0b"
        logo_url: str | None = settings.get("logo_url")

        if not prompt:
            return _response(400, {"error": "prompt is required"})
        if not api_key:
            return _response(400, {"error": "No API key saved for this user"})

        # --- Generate slide content via Qwen ---
        client = OpenAI(api_key=api_key, base_url=QWEN_BASE_URL)

        system_prompt = (
            "You are a helpful assistant that creates PowerPoint presentations. "
            "Return ONLY a JSON array where each element has exactly two keys: "
            '"title" (string) and "content" (string with bullet points separated by newlines). '
            "Do not include any other text, markdown, or explanation."
        )

        # When document file IDs are provided, prepend fileid:// references so
        # Qwen can use the uploaded documents as context for the presentation.
        if file_ids:
            file_refs = "".join(f"fileid://{fid}\n" for fid in file_ids)
            user_content = f"{file_refs}Create a presentation about: {prompt}"
        else:
            user_content = f"Create a presentation about: {prompt}"

        completion = client.chat.completions.create(
            model=QWEN_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content},
            ],
        )
        raw_content = completion.choices[0].message.content or ""
        slides_data = _parse_slides(raw_content)

        # --- Build PPTX ---
        pptx_bytes = _build_presentation(slides_data, primary_color, accent_color, logo_url)

        # --- Upload to S3 ---
        s3 = boto3.client("s3")
        key = f"presentations/{uuid.uuid4()}.pptx"
        s3.put_object(
            Bucket=S3_OUTPUT_BUCKET,
            Key=key,
            Body=pptx_bytes,
            ContentType=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )

        download_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": S3_OUTPUT_BUCKET, "Key": key},
            ExpiresIn=3600,
        )

        return _response(200, {"downloadUrl": download_url})

    except ValueError:
        return _response(401, {"error": "Unauthorized"})
    except Exception as exc:  # noqa: BLE001
        return _response(500, {"error": str(exc)})


def _response(status_code: int, body: dict) -> dict:
    return {
        "statusCode": status_code,
        "headers": {
            "Content-Type": "application/json",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Headers": "Content-Type",
            "Access-Control-Allow-Methods": "POST,OPTIONS",
        },
        "body": json.dumps(body),
    }
