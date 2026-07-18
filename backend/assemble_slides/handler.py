"""
assemble_slides Lambda: collects pre-generated python-pptx code from S3,
executes all batches sequentially against a single Presentation, adds the
references slide, uploads the final PPTX, and finalises the job.
Invoked by the last build_slides batch to complete.
"""
from __future__ import annotations

import builtins as _builtins
import io
import json
import math
import os
import re
import textwrap
from urllib import request as urlrequest

import boto3
from pptx import Presentation as _Prs
import pptx.util

# ─── CONSTANTS ────────────────────────────────────────────────────────────────

S3_OUTPUT_BUCKET = os.environ.get("OUTPUT_BUCKET", "")
SUPABASE_URL = os.environ.get("SUPABASE_URL", "").rstrip("/")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")
RESEND_API_KEY = os.environ.get("RESEND_API_KEY", "")
RESEND_FROM_EMAIL = os.environ.get("RESEND_FROM_EMAIL", "noreply@lemaiyanlabs.org")

# Must match BATCH_SIZE in build_slides/handler.py
BATCH_SIZE = 4

# ─── SUPABASE HELPERS ─────────────────────────────────────────────────────────

def _supabase_request(method: str, url: str, headers: dict, body: dict | None = None) -> object:
    payload = None
    request_headers = dict(headers)
    if body is not None:
        payload = json.dumps(body).encode("utf-8")
        request_headers["Content-Type"] = "application/json"
    req = urlrequest.Request(url=url, method=method, data=payload, headers=request_headers)
    try:
        with urlrequest.urlopen(req, timeout=10) as resp:
            raw = resp.read().decode("utf-8")
            return json.loads(raw) if raw else {}
    except Exception:  # noqa: BLE001
        return {}


def _update_job(job_id: str, **fields: object) -> None:
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        return
    url = f"{SUPABASE_URL}/rest/v1/jobs?id=eq.{job_id}"
    _supabase_request(
        "PATCH",
        url,
        {
            "apikey": SUPABASE_SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
            "Prefer": "return=minimal",
        },
        fields,  # type: ignore[arg-type]
    )


def _get_user_email(job_id: str) -> str | None:
    """Look up the user's email by joining jobs → user_settings."""
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        return None
    job_url = f"{SUPABASE_URL}/rest/v1/jobs?select=user_id&id=eq.{job_id}"
    job_data = _supabase_request("GET", job_url, {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
    })
    if not isinstance(job_data, list) or not job_data:
        return None
    user_id = (job_data[0] or {}).get("user_id")
    if not user_id:
        return None
    settings_url = f"{SUPABASE_URL}/rest/v1/user_settings?select=email&user_id=eq.{user_id}&limit=1"
    settings_data = _supabase_request("GET", settings_url, {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
    })
    if isinstance(settings_data, list) and settings_data:
        return (settings_data[0] or {}).get("email")
    return None


def _send_completion_email(
    job_id: str,
    presentation_title: str,
    download_url: str,
) -> None:
    """Send an email via the Resend API with the presentation download link."""
    if not RESEND_API_KEY:
        print(f"[{job_id}] Skipping email: RESEND_API_KEY not configured")
        return

    email = _get_user_email(job_id)
    if not email:
        print(f"[{job_id}] Skipping email: no email found for user")
        return

    payload = json.dumps({
        "from": f"PowerPoint Agent <{RESEND_FROM_EMAIL}>",
        "to": [email],
        "subject": f"Your presentation is ready: {presentation_title}",
        "html": (
            f"<p>Your presentation <strong>{presentation_title}</strong> is ready!</p>"
            f"<p><a href=\"{download_url}\">Click here to download your PPTX</a></p>"
            f"<p style=\"color:#888;font-size:12px;\">This link expires in 1 hour.</p>"
        ),
    }).encode("utf-8")

    try:
        req = urlrequest.Request(
            url="https://api.resend.com/emails",
            method="POST",
            data=payload,
            headers={
                "Authorization": f"Bearer {RESEND_API_KEY}",
                "Content-Type": "application/json",
            },
        )
        with urlrequest.urlopen(req, timeout=10) as resp:
            body = resp.read().decode("utf-8")
            print(f"[{job_id}] Email sent to {email} — Resend response: {body[:200]}")
    except Exception as exc:  # noqa: BLE001
        print(f"[{job_id}] Failed to send email to {email}: {exc}")


# ─── ICON RENDERING ──────────────────────────────────────────────────────────

def _recolor_png(png_bytes: bytes, hex_color: str) -> bytes:
    """Recolor a single-colour PNG icon to the target hex colour, preserving alpha."""
    from io import BytesIO
    from PIL import Image

    img = Image.open(BytesIO(png_bytes)).convert("RGBA")
    data = bytearray(img.tobytes())
    w, h = img.size
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    for y in range(h):
        for x in range(w):
            idx = (y * w + x) * 4
            if data[idx + 3] > 0:
                data[idx] = r
                data[idx + 1] = g
                data[idx + 2] = b
    out = BytesIO()
    Image.frombytes("RGBA", (w, h), bytes(data)).save(out, format="PNG")
    out.seek(0)
    return out.read()


# ─── EXECUTION NAMESPACE ─────────────────────────────────────────────────────

def _make_namespace(image_buffers: dict[str, bytes] | None = None) -> dict:
    """Create the restricted execution namespace with pre-injected globals."""
    import pptx
    import pptx.chart.data
    import pptx.dml.color
    import pptx.enum.chart
    import pptx.enum.text
    import pptx.enum.shapes
    import pptx.util
    from io import BytesIO
    from PIL import Image, ImageEnhance

    # ── Patch pptx to accept float coordinates ──
    import pptx.oxml.shapes.graphfrm as _graphfrm
    import pptx.oxml.simpletypes as _simpletypes

    _orig_new_graphicFrame = _graphfrm.CT_GraphicalObjectFrame.new_graphicFrame

    @classmethod
    def _int_new_graphicFrame(cls, id_: int, name: str, x: int, y: int, cx: int, cy: int):
        return _orig_new_graphicFrame.__func__(cls, id_, name, int(x), int(y), int(cx), int(cy))

    _graphfrm.CT_GraphicalObjectFrame.new_graphicFrame = _int_new_graphicFrame

    _orig_convert_to_xml = _simpletypes.BaseIntType.convert_to_xml

    @classmethod
    def _lenient_convert_to_xml(cls, value):
        if isinstance(value, float):
            value = int(value)
        return _orig_convert_to_xml.__func__(cls, value)

    _simpletypes.BaseIntType.convert_to_xml = _lenient_convert_to_xml

    _orig_validate_int = _simpletypes.BaseSimpleType.validate_int

    @staticmethod
    def _lenient_validate_int(value: int | float) -> None:
        if isinstance(value, float):
            value = int(value)
        return _orig_validate_int(value)

    _simpletypes.BaseSimpleType.validate_int = _lenient_validate_int

    def _no_shadow(shape):
        try:
            shape.shadow.inherit = False
        except (NotImplementedError, AttributeError):
            pass

    def _remove_outline(shape):
        try:
            shape.line.fill.background()
        except (NotImplementedError, AttributeError):
            pass

    def _render_bullets(slide, col_x, col_w, content_y, bullets, icons, primary_color):
        from pptx.util import Inches, Pt, Cm
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_AUTO_SIZE
        from io import BytesIO

        icon_size = Pt(22)
        icon_x = col_x + Inches(0.08)
        text_x = icon_x + icon_size + Pt(6)
        text_w = col_w - Inches(0.16) - icon_size - Pt(6)

        bullet_top = content_y + Cm(0.35)

        for b in bullets:
            icon_y = bullet_top + Pt(3)

            icon_fname = b.get("icon", "")
            if icons and icon_fname in icons:
                buf = BytesIO(icons[icon_fname])
                buf.seek(0)
                slide.shapes.add_picture(buf, icon_x, icon_y, icon_size, icon_size)
            else:
                fb = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_x, icon_y, Pt(10), Pt(10))
                fb.fill.solid()
                fb.fill.fore_color.rgb = primary_color
                _no_shadow(fb)
                _remove_outline(fb)

            txt_tb = slide.shapes.add_textbox(text_x, bullet_top, text_w, Cm(1.6))
            tf = txt_tb.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            run1 = p1.add_run()
            run1.text = b.get("title", "")
            run1.font.size = Pt(10)
            run1.font.bold = True
            run1.font.color.rgb = RGBColor(0, 0, 0)

            if b.get("description"):
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = b["description"]
                run2.font.size = Pt(8)
                run2.font.color.rgb = RGBColor(0, 0, 0)

            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            bullet_top += Inches(0.6)

        return bullet_top

    _SAFE_NAMES = (
        "abs", "all", "any", "bool", "bytes", "bytearray", "dict", "enumerate",
        "filter", "float", "globals", "hasattr", "int",
        "isinstance", "iter", "len", "list", "map", "max", "min", "next", "print", "range",
        "reversed", "round", "set", "sorted", "str", "sum", "tuple", "zip",
        "Exception", "ValueError", "TypeError", "KeyError",
        "IndexError", "AttributeError", "RuntimeError", "StopIteration",
    )
    safe_builtins = {
        name: getattr(_builtins, name)
        for name in _SAFE_NAMES
        if hasattr(_builtins, name)
    }

    ns: dict = {
        "__builtins__": safe_builtins,
        "__image_buffers__": image_buffers or {},
        "Inches": pptx.util.Inches,
        "Pt": pptx.util.Pt,
        "Emu": pptx.util.Emu,
        "Cm": pptx.util.Cm,
        "RGBColor": pptx.dml.color.RGBColor,
        "PP_ALIGN": pptx.enum.text.PP_ALIGN,
        "XL_LEGEND_POSITION": pptx.enum.chart.XL_LEGEND_POSITION,
        "ChartData": pptx.chart.data.ChartData,
        "XL_CHART_TYPE": pptx.enum.chart.XL_CHART_TYPE,
        "MSO_SHAPE": pptx.enum.shapes.MSO_SHAPE,
        "MSO_ANCHOR": pptx.enum.text.MSO_ANCHOR,
        "MSO_AUTO_SIZE": pptx.enum.text.MSO_AUTO_SIZE,
        "io": io,
        "math": math,
        "json": json,
        "urlrequest": urlrequest,
        "Image": Image,
        "ImageEnhance": ImageEnhance,
        "BytesIO": BytesIO,
        "no_shadow": _no_shadow,
        "remove_outline": _remove_outline,
        "render_bullets": _render_bullets,
    }
    if image_buffers:
        ns["get_image_buf"] = _get_image_buf_factory(image_buffers)
    return ns


def _get_image_buf_factory(image_buffers: dict[str, bytes]):
    """Return a get_image_buf(url, darken=False) closure that looks up
    pre-downloaded images, wrapping them in BytesIO and optionally darkening."""
    from io import BytesIO
    from PIL import Image, ImageEnhance

    def _get_image_buf(url: str, darken: bool = False):
        data = image_buffers.get(url)
        if data is None:
            placeholder = Image.new("RGB", (1920, 1080), (50, 50, 50))
            buf = BytesIO()
            placeholder.save(buf, format="PNG")
            buf.seek(0)
            return buf
        img = Image.open(BytesIO(data))
        if darken:
            enhancer = ImageEnhance.Brightness(img)
            img = enhancer.enhance(0.6)
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf

    return _get_image_buf


class _SafeIconDict(dict):
    """dict subclass that returns the first available icon bytes when a key is
    missing, preventing the OVAL shape fallback from triggering on hallucinated
    or mis-cased icon filenames generated by the LLM."""

    def __missing__(self, key: str) -> bytes:
        if self:
            return next(iter(self.values()))
        raise KeyError(key)


# ─── ROLLBACK ────────────────────────────────────────────────────────────────

def _rollback_slides(prs: object, original_count: int) -> int:
    """Remove slides added beyond original_count to undo a failed batch execution."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    removed = 0
    while len(sldIdLst) > original_count:
        sld_id = sldIdLst[-1]
        r_id = sld_id.get(qn('r:id'))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:  # noqa: BLE001
                pass
        sldIdLst.remove(sld_id)
        removed += 1
    return removed


# ─── REFERENCES SLIDE ─────────────────────────────────────────────────────────

def _add_references_slide(prs: object, structure: dict, settings: dict) -> None:
    """Append a References slide listing all source URLs from slide notes."""
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    import re as _re

    slides_data = structure.get("slides", [])
    citation_map: dict[str, str] = structure.get("_citation_map", {})

    def _final_resolve(text: str) -> str:
        if not citation_map:
            return text
        text = _re.sub(
            r'\[(\d+(?:\s*,\s*\d+)*)\]',
            lambda m: ' | '.join(
                citation_map.get(n.strip(), f'[{n.strip()}]')
                for n in m.group(1).split(',')
            ),
            text,
        )
        text = _re.sub(
            r'\[(\d+)\]',
            lambda m: citation_map.get(m.group(1), m.group(0)),
            text,
        )
        return _re.sub(r'(\s*\|\s*)+', ' | ', text)

    refs: list[tuple[str, str, str]] = []
    for slide in slides_data:
        notes = (slide.get("notes") or "").strip()
        notes = _final_resolve(notes)
        if not notes:
            continue
        refs.append((
            slide.get("slide_title", ""),
            (slide.get("sources") or "").strip(),
            notes,
        ))

    if not refs:
        return

    sw = prs.slide_width
    sh = prs.slide_height
    ref_slide = prs.slides.add_slide(prs.slide_layouts[6])

    ref_slide.background.fill.solid()
    ref_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    primary_hex = settings.get("primary_color", "#C00000").lstrip("#")
    primary_rgb = RGBColor(
        int(primary_hex[0:2], 16),
        int(primary_hex[2:4], 16),
        int(primary_hex[4:6], 16),
    )

    # Title
    title_tb = ref_slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Cm(1.0))
    p = title_tb.text_frame.paragraphs[0]
    p.text = "References"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = primary_rgb

    # Thin separator bar
    sep = ref_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.0), sw - Inches(0.8), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = primary_rgb
    try:
        sep.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass
    try:
        sep.line.fill.background()
    except Exception:  # noqa: BLE001
        pass

    # References content
    content_tb = ref_slide.shapes.add_textbox(Inches(0.4), Inches(1.15), sw - Inches(0.8), sh - Inches(1.45))
    ctf = content_tb.text_frame
    ctf.word_wrap = True

    first = True
    for title_str, sources_str, notes_str in refs:
        if first:
            p1 = ctf.paragraphs[0]
            first = False
        else:
            p_gap = ctf.add_paragraph()
            p_gap.font.size = Pt(4)
            p1 = ctf.add_paragraph()

        p1.text = title_str + (f"  |  {sources_str}" if sources_str else "")
        p1.font.bold = True
        p1.font.size = Pt(9)
        p1.font.color.rgb = RGBColor(40, 40, 40)

        p2 = ctf.add_paragraph()
        p2.text = notes_str
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(100, 100, 100)


# ─── MAIN HANDLER ─────────────────────────────────────────────────────────────

def handler(event: dict, context) -> None:  # noqa: ANN001
    """Entry point — invoked when all build_slides batches have saved their code.
    Loads all code files, executes them sequentially against a single Presentation,
    adds the references slide, and finalises the job."""
    job_id: str = event["job_id"]
    structure: dict = event["structure"]
    settings: dict = event.get("settings", {})
    total_batches: int = event["total_batches"]

    print(f"[{job_id}] Assemble-slides started. {total_batches} batches expected.")

    s3 = boto3.client("s3")

    # ── Idempotency lock ───────────────────────────────────────────────────
    lock_key = f"wip/{job_id}/assembling"
    try:
        s3.head_object(Bucket=S3_OUTPUT_BUCKET, Key=lock_key)
        print(f"[{job_id}] Assembly lock already exists — another invocation is handling this job. Exiting.")
        return
    except Exception:  # noqa: BLE001
        pass
    s3.put_object(Bucket=S3_OUTPUT_BUCKET, Key=lock_key, Body=b"")

    try:
        # ── Load all batch code files ──────────────────────────────────────
        code_blocks: dict[int, str] = {}
        for i in range(1, total_batches + 1):
            try:
                resp = s3.get_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/code_{i}.txt")
                code_blocks[i] = resp["Body"].read().decode()
                print(f"[{job_id}] Loaded code batch {i}/{total_batches} ({len(code_blocks[i])} chars)")
            except Exception:  # noqa: BLE001
                print(f"[{job_id}] WARNING: code batch {i} not found — skipping")

        if not code_blocks:
            raise RuntimeError("No code batches found — all batches may have failed")

        # ── Download logo ──────────────────────────────────────────────────
        logo_url: str = settings.get("logo_url", "")
        logo_bytes: bytes | None = None
        if logo_url:
            try:
                req = urlrequest.Request(url=logo_url, method="GET")
                with urlrequest.urlopen(req, timeout=15) as resp:
                    logo_bytes = resp.read()
                print(f"[{job_id}] Logo downloaded ({len(logo_bytes)} bytes)")
            except Exception:  # noqa: BLE001
                logo_bytes = None

        # ── Pre-download all background images ─────────────────────────────
        slides: list = structure.get("slides", [])
        image_urls: set[str] = set()
        for slide in slides:
            url = slide.get("image_url")
            if url:
                image_urls.add(url)
        image_buffers: dict[str, bytes] = {}
        if image_urls:
            print(f"[{job_id}] Downloading {len(image_urls)} unique background images...")
            for url in sorted(image_urls):
                try:
                    req = urlrequest.Request(url=url, method="GET")
                    with urlrequest.urlopen(req, timeout=15) as resp:
                        image_buffers[url] = resp.read()
                except Exception:  # noqa: BLE001
                    print(f"[{job_id}] Failed to download image {url[:80]}..., skipping")
            print(f"[{job_id}] Downloaded {len(image_buffers)}/{len(image_urls)} images")

        # ── Download all icons ─────────────────────────────────────────────
        icons: dict[str, bytes] = {}
        icon_names: set[str] = set()
        for slide in slides:
            for col in slide.get("columns", []):
                for bullet in col.get("bullets") or []:
                    icon_name = bullet.get("icon", "")
                    if icon_name:
                        icon_names.add(icon_name)
        if icon_names:
            primary_color = settings.get("primary_color", "#C00000")
            print(f"[{job_id}] Downloading {len(icon_names)} unique icons...")
            for name in sorted(icon_names):
                try:
                    buf = io.BytesIO()
                    png_name = name.replace(".svg", ".png")
                    s3.download_fileobj(S3_OUTPUT_BUCKET, f"icons/{png_name}", buf)
                    recolored = _recolor_png(buf.getvalue(), primary_color)
                    icons[name] = recolored
                except Exception:  # noqa: BLE001
                    print(f"[{job_id}] Failed to download icon {name}, skipping")
            print(f"[{job_id}] Downloaded {len(icons)}/{len(icon_names)} icons")

        has_logo = logo_bytes is not None
        has_icons = len(icons) > 0

        # ── Create Presentation & execute all batches sequentially ─────────
        prs = _Prs()
        prs.slide_width = pptx.util.Inches(13.33)
        prs.slide_height = pptx.util.Inches(7.5)

        namespace = _make_namespace(image_buffers)
        total_slides = len(slides)

        for batch_num in sorted(code_blocks):
            code = code_blocks[batch_num]
            fn_name = f"build_batch_{batch_num}"
            start_slide = (batch_num - 1) * BATCH_SIZE + 1
            end_slide = min(batch_num * BATCH_SIZE, total_slides)
            batch_slides_data = slides[(batch_num - 1) * BATCH_SIZE : batch_num * BATCH_SIZE]

            _update_job(job_id, stage_message=f"Assembling slides {start_slide}-{end_slide} of {total_slides}…")

            # Execute with up to 3 retries (in case of code errors)
            batch_ok = False
            last_error = ""
            for attempt in range(3):
                slide_count_before = len(prs.slides._sldIdLst)
                try:
                    code_clean = re.sub(r"```(?:python)?\s*", "", code).strip().rstrip("`").strip()
                    exec(code_clean, namespace)  # noqa: S102

                    build_fn = namespace.get(fn_name)
                    if not callable(build_fn):
                        raise ValueError(f"Generated code does not define '{fn_name}'")

                    kwargs = {
                        "logo_bytes": logo_bytes,
                        "icons": _SafeIconDict(icons or {}),
                    }
                    import inspect as _inspect
                    sig = _inspect.signature(build_fn)
                    if "image_buffers" in sig.parameters:
                        kwargs["image_buffers"] = namespace.get("__image_buffers__", {})
                    if "slides_data" in sig.parameters:
                        kwargs["slides_data"] = batch_slides_data
                    build_fn(prs, **kwargs)
                    print(f"[{job_id}] Batch {batch_num}/{total_batches}: executed on attempt {attempt + 1}")
                    batch_ok = True
                    break
                except Exception as exc:  # noqa: BLE001
                    _rollback_slides(prs, slide_count_before)
                    print(f"[{job_id}] Batch {batch_num}/{total_batches}: execution attempt {attempt + 1} FAILED: {exc}")
                    last_error = str(exc)

            if not batch_ok:
                print(f"[{job_id}] WARNING: Batch {batch_num}/{total_batches} failed after 3 attempts: {last_error}")
                # Continue with remaining batches — partial presentation is better than nothing

        # ── Add references slide ────────────────────────────────────────────
        print(f"[{job_id}] Adding references slide...")
        _add_references_slide(prs, structure, settings)

        # ── Aggregate all code for Supabase ────────────────────────────────
        all_code: list[str] = [code_blocks[i] for i in sorted(code_blocks)]
        _update_job(job_id, pptx_code="\n\n# --- BATCH ---\n\n".join(all_code))

        # ── Upload final PPTX ──────────────────────────────────────────────
        prs_buf = io.BytesIO()
        prs.save(prs_buf)
        prs_bytes = prs_buf.getvalue()

        key = f"presentations/{job_id}.pptx"
        s3.put_object(
            Bucket=S3_OUTPUT_BUCKET,
            Key=key,
            Body=prs_bytes,
            ContentType=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )
        download_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": S3_OUTPUT_BUCKET, "Key": key},
            ExpiresIn=259200,
        )
        _update_job(
            job_id,
            status="done",
            stage_message="Your presentation is ready!",
            download_url=download_url,
        )
        print(f"[{job_id}] DONE. Download: {download_url}")

        # ── Send completion email ───────────────────────────────────────────
        title = structure.get("presentation_title", "Your Presentation")
        _send_completion_email(job_id, title, download_url)

        # ── Clean up WIP files ──────────────────────────────────────────────
        print(f"[{job_id}] Cleaning up WIP files...")
        for i in range(1, total_batches + 1):
            try:
                s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/code_{i}.txt")
            except Exception:  # noqa: BLE001
                pass
            try:
                s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/done_{i}")
            except Exception:  # noqa: BLE001
                pass
            try:
                s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=f"wip/{job_id}/failed_{i}")
            except Exception:  # noqa: BLE001
                pass
        try:
            s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=lock_key)
        except Exception:  # noqa: BLE001
            pass

    except Exception as exc:  # noqa: BLE001
        print(f"[{job_id}] ERROR: {exc}")
        _update_job(
            job_id,
            status="error",
            stage_message="Something went wrong during assembly.",
            error_message=str(exc),
        )
        # Remove lock so a retry can proceed
        try:
            s3.delete_object(Bucket=S3_OUTPUT_BUCKET, Key=lock_key)
        except Exception:  # noqa: BLE001
            pass
        raise
